from django.shortcuts import render
from django.http import JsonResponse, HttpResponse, Http404
from django.views.decorators.csrf import csrf_exempt
from django.conf import settings
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches

def index(request):
    """Page principale avec l'interface de upload"""
    return render(request, 'file_processor/index.html')

@csrf_exempt
def upload_file(request):
    """Traite le fichier Excel uploadé et génère les PowerPoints"""
    if request.method != 'POST':
        return JsonResponse({'error': 'Méthode non autorisée'}, status=405)

    if 'file' not in request.FILES:
        return JsonResponse({'error': 'Aucun fichier fourni'}, status=400)

    file = request.FILES['file']
    sort_column = request.POST.get('sort_column', 'code site')

    # Récupérer les informations de visite
    date_debut = request.POST.get('date_debut', '')
    date_fin = request.POST.get('date_fin', '')
    objet_visite = request.POST.get('objet_visite', '')

    # Vérifier l'extension du fichier
    if not file.name.endswith(('.xlsx', '.xls')):
        return JsonResponse({'error': 'Le fichier doit être un fichier Excel (.xlsx ou .xls)'}, status=400)

    try:
        # Calculer la taille du fichier
        file.seek(0, 2)  # Aller à la fin du fichier
        file_size_bytes = file.tell()
        file.seek(0)  # Revenir au début

        # Convertir en format lisible
        if file_size_bytes < 1024:
            file_size_display = f"{file_size_bytes} B"
        elif file_size_bytes < 1024 * 1024:
            file_size_display = f"{file_size_bytes / 1024:.1f} KB"
        else:
            file_size_display = f"{file_size_bytes / (1024 * 1024):.1f} MB"

        # Lire le fichier Excel
        df = pd.read_excel(file)
        print('Colonnes Excel:', list(df.columns))

        # Normaliser les noms de colonnes (supprimer espaces en début/fin et convertir en minuscules pour comparaison)
        df.columns = df.columns.str.strip()

        # Créer un mapping flexible pour les noms de colonnes
        column_mapping = {}

        # Mapping des colonnes avec recherche très flexible
        def normalize_column_name(name):
            """Normalise un nom de colonne pour la comparaison"""
            return ''.join(name.lower().split())  # Supprime tous les espaces et met en minuscules

        column_searches = {
            'code site': ['codesite'],
            'ST FO': ['stfo'],
            'contact ERPT': ['contacterpt'],
            'Contact IAM': ['contactiam', 'contact iam', 'contact_iam'],
            'DR IAM': ['driam'],
            'ville': ['ville'],
            'Date TSS': ['datetss', 'date tss', 'date_tss'],
            'X Départ ERPT - Y Départ ERPT': ['xdéparterpt-ydéparterpt', 'xdéparterptydepart'],
            'X Arrivée ERPT Proposition1 - Y Arrivée ERPT Proposition1': ['xarrivéeerptproposition1-yarrivéeerptproposition1', 'xarriveeerpt-yarriveeerpt']
        }

        # Recherche très flexible des colonnes
        for target_col, search_terms in column_searches.items():
            found = False
            for actual_col in df.columns:
                actual_col_normalized = normalize_column_name(actual_col)
                for search_term in search_terms:
                    if search_term in actual_col_normalized or actual_col_normalized in search_term:
                        column_mapping[target_col] = actual_col
                        found = True
                        break
                if found:
                    break

        # Vérifier les colonnes essentielles
        essential_columns = ['code site', 'ST FO', 'DR IAM', 'ville']
        missing_columns = [col for col in essential_columns if col not in column_mapping]

        if missing_columns:
            return JsonResponse({
                'error': f'Colonnes essentielles manquantes: {", ".join(missing_columns)}. Colonnes disponibles: {", ".join(df.columns)}',
                'available_columns': list(df.columns),
                'column_mapping': column_mapping
            }, status=400)

        # Renommer les colonnes pour correspondre aux noms attendus
        rename_dict = {v: k for k, v in column_mapping.items()}
        df = df.rename(columns=rename_dict)

        # Trier les données selon la colonne choisie
        if sort_column not in df.columns:
            return JsonResponse({'error': f'Colonne de tri "{sort_column}" non trouvée'}, status=400)

        # Calculer les statistiques selon la colonne de tri
        total_rows = len(df)

        # Statistiques par colonne de tri
        sort_stats = df[sort_column].value_counts().to_dict()

        # Adapter les statistiques selon la colonne choisie avec détails enrichis
        if sort_column == 'DR IAM':
            # Pour DR IAM : afficher DR stats avec ST FO et fichiers
            primary_stats = {}
            for dr_name in df['DR IAM'].unique():
                dr_data = df[df['DR IAM'] == dr_name]
                lines_count = len(dr_data)
                st_fo_count = len(dr_data['ST FO'].unique())
                # Estimer le nombre de fichiers (basé sur 19 lignes max par fichier)
                files_count = max(1, (lines_count + 18) // 19)  # Arrondi vers le haut
                primary_stats[dr_name] = {
                    'lines': lines_count,
                    'st_fo': st_fo_count,
                    'files': files_count
                }
            primary_count = len(primary_stats)
            primary_label = 'DR'
            stats_title = 'Répartition par DR IAM'
            show_details = True
        elif sort_column == 'ville':
            # Pour ville : afficher ville stats avec ST FO et fichiers
            primary_stats = {}
            for ville_name in df['ville'].unique():
                ville_data = df[df['ville'] == ville_name]
                lines_count = len(ville_data)
                st_fo_count = len(ville_data['ST FO'].unique())
                # Estimer le nombre de fichiers (basé sur 19 lignes max par fichier)
                files_count = max(1, (lines_count + 18) // 19)  # Arrondi vers le haut
                primary_stats[ville_name] = {
                    'lines': lines_count,
                    'st_fo': st_fo_count,
                    'files': files_count
                }
            primary_count = len(primary_stats)
            primary_label = 'ville'
            stats_title = 'Répartition par ville'
            show_details = True
        elif sort_column == 'ST FO':
            # Pour ST FO : afficher ST FO stats avec lignes et fichiers
            primary_stats = {}
            for st_fo_name in df['ST FO'].unique():
                st_fo_data = df[df['ST FO'] == st_fo_name]
                lines_count = len(st_fo_data)
                # Pour ST FO, le nombre de ST FO est toujours 1
                st_fo_count = 1
                # Estimer le nombre de fichiers (basé sur 19 lignes max par fichier)
                files_count = max(1, (lines_count + 18) // 19)  # Arrondi vers le haut
                primary_stats[st_fo_name] = {
                    'lines': lines_count,
                    'st_fo': st_fo_count,
                    'files': files_count
                }
            primary_count = len(primary_stats)
            primary_label = 'ST FO'
            stats_title = 'Répartition par ST FO'
            show_details = True
        else:  # code site - simple (pas de détails)
            # Pour code site : pas de statistiques détaillées
            primary_stats = {}
            primary_count = len(df['code site'].value_counts())
            primary_label = 'code site'
            stats_title = ''
            show_details = False

        # Éviter la duplication des statistiques pour DR IAM, ville et ST FO
        if sort_column in ['DR IAM', 'ville', 'ST FO']:
            # Ne pas envoyer sort_stats pour éviter la duplication
            final_sort_stats = {}
        else:
            final_sort_stats = sort_stats

        # Calculer le total des ST FO uniques dans le fichier
        total_st_fo = len(df['ST FO'].unique())

        # Statistiques générales
        stats = {
            'total_rows': total_rows,
            'total_st_fo': total_st_fo,
            'primary_count': primary_count,
            'primary_stats': primary_stats,
            'primary_label': primary_label,
            'stats_title': stats_title,
            'sort_column': sort_column,
            'sort_stats': final_sort_stats,
            'group_count': len(final_sort_stats),
            'show_details': show_details,
            'file_size': file_size_display
        }

        # Vérification des colonnes juste avant le groupby
        print('Colonnes du DataFrame juste avant groupby:', list(df.columns))
        print('Colonne de tri demandée:', sort_column)
        # Grouper les données par la colonne de tri
        grouped_data = df.groupby(sort_column)
        print('Groupes trouvés pour le tri:', list(grouped_data.groups.keys()))

        # Créer un dossier temporaire pour les fichiers PowerPoint
        output_dir = os.path.join(settings.MEDIA_ROOT, 'generated_ppts')
        os.makedirs(output_dir, exist_ok=True)

        # Nettoyer les anciens fichiers du même type de tri
        sort_prefix = sort_column.replace(' ', ' ')  # Garder les espaces comme dans les noms de fichiers
        cleaned_count = 0
        for existing_file in os.listdir(output_dir):
            if existing_file.startswith(f"{sort_prefix}_") and existing_file.endswith('.pptx'):
                os.remove(os.path.join(output_dir, existing_file))
                cleaned_count += 1


        generated_files = []
        created_files = []  # Correction pour éviter l'erreur de variable non définie
        # Générer un PowerPoint pour chaque groupe (ordre d'origine)
        for group_name, group_df in grouped_data:
            print('Génération pour le groupe:', group_name, 'lignes:', len(group_df))
            ppt_filename = f"{sort_column}_{group_name}.pptx"
            ppt_path = os.path.join(output_dir, ppt_filename)
            created_files.extend(create_powerpoint(group_df, ppt_path, group_name, sort_column, date_debut, date_fin, objet_visite))

        # Après la génération des fichiers
        # Diagnostic : afficher le contenu de created_files
        print('created_files pour le calcul du total :', created_files)
        # Afficher le nombre total de fichiers PowerPoint générés
        stats['file_count'] = len(created_files)
        return JsonResponse({
            'success': True,
            'message': f'{len(created_files)} fichiers PowerPoint générés',
            'files': [file_info['filename'] for file_info in created_files],
            'files_details': created_files,  # Clé essentielle pour l'affichage côté frontend
            'stats': stats
        })

    except Exception as e:
        return JsonResponse({'error': f'Erreur lors du traitement: {str(e)}'}, status=500)



def create_powerpoint(df, output_path, group_name, sort_column, date_debut='', date_fin='', objet_visite=''):
    """Crée des présentations PowerPoint à partir des données avec division en fichiers (max 19 lignes par fichier)"""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt, Cm


    # Colonnes finales à afficher dans le PowerPoint
    final_columns = [
        'code site',
        'ST FO',
        'contact ERPT',
        'Contact IAM',
        'DR IAM',
        'ville',
        'Date TSS',
        'X Départ ERPT - Y Départ ERPT',
        'X Arrivée ERPT Proposition1 - Y Arrivée ERPT Proposition1'
    ]
    display_columns = [
        'code site',
        'ST FO',
        'contact ERPT',
        'Contact IAM',
        'DR IAM',
        'ville',
        'Date TSS',
        'X Départ ERPT - Y Départ ERPT',
        'X Arrivée ERPT - Y Arrivée ERPT'
    ]

    # Chemin vers l'image AAA
    image_path = 'AAA.jpeg'
    image_exists = os.path.exists(image_path)

    # Créer une copie du DataFrame pour éviter de modifier l'original
    df_work = df.copy().reset_index(drop=True)  # Reset index pour éviter les problèmes d'ordre

    # Créer les colonnes manquantes avec des valeurs par défaut SEULEMENT si elles n'existent pas
    if 'contact ERPT' not in df_work.columns:
        df_work['contact ERPT'] = [f"contact{i+1}@erpt.fr" for i in range(len(df_work))]
    if 'Contact IAM' not in df_work.columns:
        df_work['Contact IAM'] = ["" for _ in range(len(df_work))]
    if 'Date TSS' not in df_work.columns:
        df_work['Date TSS'] = ["" for _ in range(len(df_work))]

    # Gérer les colonnes X-Y (soit séparées, soit déjà combinées)
    # Vérifier si les colonnes sont déjà combinées dans le fichier Excel
    if 'X Départ ERPT - Y Départ ERPT' not in df_work.columns:
        # Les colonnes ne sont pas encore combinées, essayer de les combiner
        if 'X Départ ERPT' in df_work.columns and 'Y Départ ERPT' in df_work.columns:
            combined_depart = []
            for _, row in df_work.iterrows():
                x_val = row['X Départ ERPT']
                y_val = row['Y Départ ERPT']
                if pd.notna(x_val) and pd.notna(y_val) and str(x_val).strip() != '' and str(y_val).strip() != '':
                    combined_depart.append(f"{x_val} - {y_val}")
                else:
                    combined_depart.append("")
            df_work['X Départ ERPT - Y Départ ERPT'] = combined_depart
        else:
            df_work['X Départ ERPT - Y Départ ERPT'] = [""] * len(df_work)
    # Si la colonne existe déjà, s'assurer qu'elle contient des données valides
    else:
        # Nettoyer les valeurs NaN et vides
        df_work['X Départ ERPT - Y Départ ERPT'] = df_work['X Départ ERPT - Y Départ ERPT'].fillna("").astype(str)

    if 'X Arrivée ERPT - Y Arrivée ERPT' not in df_work.columns:
        # Les colonnes ne sont pas encore combinées, essayer de les combiner
        if 'X Arrivée ERPT Proposition1' in df_work.columns and 'Y Arrivée ERPT Proposition1' in df_work.columns:
            combined_arrivee = []
            for _, row in df_work.iterrows():
                x_val = row['X Arrivée ERPT Proposition1']
                y_val = row['Y Arrivée ERPT Proposition1']
                if pd.notna(x_val) and pd.notna(y_val) and str(x_val).strip() != '' and str(y_val).strip() != '':
                    combined_arrivee.append(f"{x_val} - {y_val}")
                else:
                    combined_arrivee.append("")
            df_work['X Arrivée ERPT - Y Arrivée ERPT'] = combined_arrivee
        else:
            df_work['X Arrivée ERPT - Y Arrivée ERPT'] = [""] * len(df_work)
    # Si la colonne existe déjà, s'assurer qu'elle contient des données valides
    else:
        # Nettoyer les valeurs NaN et vides
        df_work['X Arrivée ERPT - Y Arrivée ERPT'] = df_work['X Arrivée ERPT - Y Arrivée ERPT'].fillna("").astype(str)

    # Créer les autres colonnes manquantes si nécessaire
    for col in final_columns:
        if col not in df_work.columns:
            df_work[col] = ["" for _ in range(len(df_work))]

    # Sélectionner seulement les colonnes finales à afficher dans l'ordre exact
    df_filtered = df_work[final_columns].copy()

    # Nettoyer les données
    for col in df_filtered.columns:
        df_filtered[col] = df_filtered[col].fillna("").astype(str)

    # Appliquer le deuxième tri par ST FO pour diviser en sous-groupes
    df_filtered = df_filtered.sort_values('ST FO', na_position='last').reset_index(drop=True)

    # Grouper par ST FO (deuxième niveau de tri)
    st_fo_groups = df_filtered.groupby('ST FO', sort=False)

    # Optimiser la distribution des ST FO dans les chunks
    chunks = []
    chunk_names = []

    # Collecter tous les groupes ST FO avec leurs tailles
    st_fo_data_list = []
    for st_fo_name, st_fo_group in st_fo_groups:
        st_fo_data_list.append({
            'name': st_fo_name,
            'data': st_fo_group.reset_index(drop=True),
            'size': len(st_fo_group)
        })

    # Algorithme d'optimisation pour remplir les chunks
    current_chunk = pd.DataFrame()
    current_chunk_st_fo = []
    remaining_st_fo = st_fo_data_list.copy()

    while remaining_st_fo:
        current_space = 19 - len(current_chunk)

        if current_space <= 0:
            # Chunk plein, le sauvegarder
            chunks.append(current_chunk.copy())
            if len(current_chunk_st_fo) == 1:
                chunk_names.append(current_chunk_st_fo[0])
            else:
                chunk_names.append(" + ".join(current_chunk_st_fo))

            # Commencer un nouveau chunk
            current_chunk = pd.DataFrame()
            current_chunk_st_fo = []
            continue

        # Chercher le meilleur ST FO à ajouter
        best_fit = None
        best_index = -1

        for i, st_fo_item in enumerate(remaining_st_fo):
            if st_fo_item['size'] <= current_space:
                # Ce ST FO peut entrer entièrement
                if best_fit is None or st_fo_item['size'] > best_fit['size']:
                    best_fit = st_fo_item
                    best_index = i

        if best_fit:
            # Ajouter le ST FO entier au chunk
            if len(current_chunk) == 0:
                current_chunk = best_fit['data'].copy()
            else:
                current_chunk = pd.concat([current_chunk, best_fit['data']], ignore_index=True)
            current_chunk_st_fo.append(best_fit['name'])
            remaining_st_fo.pop(best_index)
        else:
            # Aucun ST FO ne peut entrer entièrement, prendre une partie du plus grand
            if remaining_st_fo:
                largest_st_fo = max(remaining_st_fo, key=lambda x: x['size'])
                largest_index = remaining_st_fo.index(largest_st_fo)

                # Prendre autant de lignes que possible
                lines_to_take = min(current_space, largest_st_fo['size'])
                partial_data = largest_st_fo['data'].iloc[:lines_to_take].copy()

                if len(current_chunk) == 0:
                    current_chunk = partial_data
                else:
                    current_chunk = pd.concat([current_chunk, partial_data], ignore_index=True)

                # Mettre à jour le nom du chunk
                if lines_to_take == largest_st_fo['size']:
                    # ST FO entier utilisé
                    current_chunk_st_fo.append(largest_st_fo['name'])
                    remaining_st_fo.pop(largest_index)
                else:
                    # ST FO partiellement utilisé
                    current_chunk_st_fo.append(f"{largest_st_fo['name']} (partiel)")
                    # Mettre à jour le ST FO restant
                    remaining_st_fo[largest_index]['data'] = largest_st_fo['data'].iloc[lines_to_take:].reset_index(drop=True)
                    remaining_st_fo[largest_index]['size'] = len(remaining_st_fo[largest_index]['data'])

    # Ajouter le dernier chunk s'il n'est pas vide
    if len(current_chunk) > 0:
        chunks.append(current_chunk)
        if len(current_chunk_st_fo) == 1:
            chunk_names.append(current_chunk_st_fo[0])
        else:
            chunk_names.append(" + ".join(current_chunk_st_fo))

    # Créer plusieurs fichiers PowerPoint si nécessaire
    created_files = []

    for chunk_index, chunk_data in enumerate(chunks):
        # Créer un nouveau PowerPoint pour chaque chunk
        prs = Presentation()

        # Déterminer le nom du fichier basé sur ST FO
        chunk_st_fo_name = chunk_names[chunk_index]
        chunk_lines = len(chunk_data)

        # Slide de titre avec modifications CSS simples
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"Données pour {sort_column}: {group_name} - {chunk_st_fo_name}"
        subtitle.text = f"Lignes: {chunk_lines} | ST FO: {chunk_st_fo_name}"

        # Appliquer la police Arial Narrow au titre avec style amélioré
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.name = 'Arial Narrow'
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Noir

        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.name = 'Arial Narrow'
        subtitle_paragraph.font.size = Pt(18)
        subtitle_paragraph.font.color.rgb = RGBColor(64, 64, 64)  # Gris foncé

        # Ajouter une bordure noire autour du slide (diminuée de 4cm largeur et 4cm hauteur)
        try:
            # Créer un rectangle de bordure réduit
            from pptx.enum.shapes import MSO_SHAPE
            left = Cm(2.2)    # 2cm de plus de chaque côté (0.2 + 2.0)
            top = Cm(2.2)     # 2cm de plus de chaque côté (0.2 + 2.0)
            width = Cm(21.2)  # Largeur réduite de 4cm (25.2 - 4.0)
            height = Cm(14.8) # Hauteur réduite de 4cm (18.8 - 4.0)

            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )

            # Style de la bordure
            border_fill = border_shape.fill
            border_fill.background()  # Fond transparent

            border_line = border_shape.line
            border_line.color.rgb = RGBColor(0, 0, 0)  # Bordure noire
            border_line.width = Pt(3)  # Épaisseur 3pt

            # Envoyer la bordure à l'arrière-plan
            border_shape.element.getparent().remove(border_shape.element)
            slide.shapes._spTree.insert(2, border_shape.element)

        except Exception as e:
            # Si erreur avec la bordure, continuer sans
            pass
        # Slide avec l'image AAA et le tableau en dessous (sans titre)
        blank_slide_layout = prs.slide_layouts[6]  # Layout vide sans titre
        data_slide = prs.slides.add_slide(blank_slide_layout)

        # Ajouter l'image AAA si elle existe
        if image_exists:
            try:
                # Positionner l'image tout en haut de la slide (sans titre)
                left = Cm(1.0)
                top = Cm(0.5)  # Tout en haut avec petite marge
                width = Cm(24.5)  # Largeur diminuée de 1.5cm (26.0 - 1.5 = 24.5)
                height = Cm(7.0)  # Hauteur réduite pour laisser plus de place au tableau

                data_slide.shapes.add_picture(image_path, left, top, width, height)
            except Exception as e:
                # Si erreur avec l'image, continuer sans l'image
                pass

        # Ajouter les informations de visite EN AVANT-PLAN sur l'image AAA
        if date_debut or date_fin or objet_visite:

            # Formater les dates au format DD/MM/YYYY
            def format_date(date_str):
                if date_str:
                    try:
                        from datetime import datetime
                        # Convertir de YYYY-MM-DD vers DD/MM/YYYY
                        date_obj = datetime.strptime(date_str, '%Y-%m-%d')
                        return date_obj.strftime('%d/%m/%Y')
                    except:
                        return date_str
                return ',,/,,/,,,,'

            date_debut_formatted = format_date(date_debut)
            date_fin_formatted = format_date(date_fin)
            objet_formatted = f'"{objet_visite}"' if objet_visite else '""'

            # Créer le texte dans le format exact demandé
            visit_info_text = f"{date_debut_formatted}         {date_fin_formatted}                                           {objet_formatted}"

            # Ajouter la zone de texte EN AVANT-PLAN sur l'image AAA
            # Position ajustée : descendre de 0.2cm supplémentaire
            text_left = Cm(4.7)   # Pousser à gauche 0.3cm (inchangé)
            text_top = Cm(6.98)   # Descendre de 0.2cm (6.78 + 0.2 = 6.98)
            text_width = Cm(24.0) # Largeur adaptée à l'image
            text_height = Cm(1.2) # Hauteur suffisante pour le texte

            try:
                text_box = data_slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
                text_frame = text_box.text_frame
                text_frame.text = visit_info_text

                # Style du texte pour visibilité sur l'image
                paragraph = text_frame.paragraphs[0]
                paragraph.font.name = 'Arial Narrow'
                paragraph.font.size = Pt(9)  # Taille réduite à 9pt
                paragraph.font.bold = False  # Pas de gras
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Texte noir

                # Fond TRANSPARENT (pas de fond)
                fill = text_box.fill
                fill.background()  # Fond transparent

            except Exception as e:
                # Si erreur avec la zone de texte, continuer sans
                pass

        # Ajouter le tableau en dessous de l'image pour ce chunk
        rows = len(chunk_data) + 1  # +1 pour l'en-tête
        cols = len(final_columns)

        # Calculer la largeur totale nécessaire (somme des largeurs de colonnes)
        total_width_cm = 1.8 + 2.3 + 3.4 + 3.4 + 2.5 + 2.5 + 2.4 + 3.0 + 3.0  # = 28.0 cm

        # Positionner le tableau DIRECTEMENT en dessous de l'image
        # (les informations de visite sont maintenant en avant-plan sur l'image)
        left = Cm(1.0)  # Marge à gauche
        if image_exists:
            top = Cm(7.5)  # DIRECTEMENT collé à l'image
        else:
            top = Cm(1.0)   # Plus haut si pas d'image
        width = Cm(total_width_cm)  # Largeur exacte du tableau
        height = Inches(6)

        table = data_slide.shapes.add_table(rows, cols, left, top, width, height).table

        # En-têtes avec police Arial Narrow
        for i, column in enumerate(display_columns):
            cell = table.cell(0, i)
            cell.text = str(column)

            # Appliquer la police Arial Narrow et le style aux en-têtes
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = 'Arial Narrow'
            paragraph.font.size = Pt(7)  # Taille réduite à 7pt
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Blanc

            # Réduire l'espacement dans la cellule
            cell.text_frame.margin_top = Inches(0.02)
            cell.text_frame.margin_bottom = Inches(0.02)
            cell.text_frame.margin_left = Inches(0.05)
            cell.text_frame.margin_right = Inches(0.05)

            # Couleur de fond pour l'en-tête
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0x68, 0x00)  # Orange vif (ff6800)

        # Données avec police Arial Narrow (utiliser chunk_data au lieu de df_filtered)
        for i, (_, row) in enumerate(chunk_data.iterrows(), 1):
            for j, value in enumerate(row):
                # Format spécial pour les deux dernières colonnes (coordonnées)
                if final_columns[j] in [
                    'X Départ ERPT - Y Départ ERPT',
                    'X Arrivée ERPT Proposition1 - Y Arrivée ERPT Proposition1'
                ]:
                    try:
                        if isinstance(value, str):
                            # On accepte les séparateurs virgule, espace ou tiret
                            for sep in [',', ' - ', ' ']:
                                if sep in value:
                                    parts = [p.strip() for p in value.split(sep)]
                                    break
                            else:
                                parts = [value.strip()]
                            formatted_parts = []
                            for part in parts:
                                if part:
                                    try:
                                        formatted_parts.append(f"{float(part):.5f}")
                                    except Exception:
                                        formatted_parts.append(part)
                            value = ', '.join(formatted_parts)
                        elif value != "":
                            value = f"{float(value):.5f}"
                    except Exception:
                        pass
                cell = table.cell(i, j)
                cell.text = str(value) if pd.notna(value) and str(value) != "" else ""

                # Appliquer la police Arial Narrow aux données
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.name = 'Arial Narrow'
                paragraph.font.size = Pt(7)  # Taille réduite à 7pt

                # Réduire l'espacement dans la cellule de données
                cell.text_frame.margin_top = Inches(0.02)
                cell.text_frame.margin_bottom = Inches(0.02)
                cell.text_frame.margin_left = Inches(0.05)
                cell.text_frame.margin_right = Inches(0.05)

                # Couleur alternée pour les lignes
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xD4, 0xB9)  # Orange clair (ffd4b9)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # Blanc (ffffff)

        # Largeurs de colonnes spécifiées en centimètres selon la nouvelle demande
        from pptx.util import Cm
        column_widths = [
            Cm(1.7),   # colonne 1: code site - 1,7 cm
            Cm(2.4),   # colonne 2: ST FO - 2,4 cm
            Cm(3.1),   # colonne 3: contact ERPT - 3,1 cm
            Cm(3.1),   # colonne 4: Contact IAM - 3,1 cm
            Cm(2.5),   # colonne 5: DR IAM - 2,5 cm
            Cm(2.5),   # colonne 6: ville - 2,5 cm
            Cm(2.3),   # colonne 7: Date TSS - 2,3 cm
            Cm(2.8),   # colonne 8: X Départ ERPT - Y Départ ERPT - 2,8 cm
            Cm(2.8)    # colonne 9: X Arrivée ERPT Proposition1 - Y Arrivée - 2,8 cm
        ]
        # Appliquer les largeurs aux colonnes
        for i, width in enumerate(column_widths):
            if i < len(table.columns):
                table.columns[i].width = width

        # Ajuster la hauteur des lignes - format uniforme ultra-compact

        for i, row in enumerate(table.rows):
            if i == 0:  # En-tête
                row.height = Cm(0.10)  # En-tête: 0.10 cm
            else:  # Toutes les lignes de données
                row.height = Cm(0.05)  # Toutes les données: 0.05 cm



        # Sauvegarder ce fichier PowerPoint avec nom basé sur ST FO
        base_name = output_path.replace('.pptx', '')
        # Nettoyer le nom ST FO pour le nom de fichier (enlever caractères spéciaux)
        safe_st_fo_name = chunk_st_fo_name.replace('/', '_').replace('\\', '_').replace(':', '_')
        chunk_output_path = f"{base_name} {safe_st_fo_name}.pptx"

        prs.save(chunk_output_path)
        # Ajouter les informations détaillées du fichier
        created_files.append({
            'filename': os.path.basename(chunk_output_path),
            'lines': len(chunk_data)
        })

    # Retourner la liste des fichiers créés avec détails
    print('Fichiers créés pour le groupe', group_name, ':', created_files)
    return created_files

def download_file(request, filename):
    """Permet de télécharger un fichier PowerPoint généré"""
    file_path = os.path.join(settings.MEDIA_ROOT, 'generated_ppts', filename)

    if not os.path.exists(file_path):
        raise Http404("Fichier non trouvé")

    with open(file_path, 'rb') as f:
        response = HttpResponse(f.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response
